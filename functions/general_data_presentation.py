import numpy as np
from pptx import Presentation
from pptx.util import Inches
import io


def smooth(y, box_pts):
    box = np.ones(box_pts)/box_pts
    y_smooth = np.convolve(y, box, mode='same')
    return y_smooth


def open_presentation():
    prs = Presentation()
    return prs


def add_txt_to_slide_ppt(slide, slide_title, left, top, width, height):
    txt_box = slide.shapes.add_textbox(left, top, width, height)
    txt_frame = txt_box.text_frame
    txt_frame.text = slide_title


def add_figure_slide_to_ppt(prs, plt, left, top, width, height):
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    graph_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(graph_slide_layout)
    # placeholder = slide.placeholders[1]
    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    return slide


def add_figure_to_existing_slide(slide,  plt, left, top, width, height):
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    return slide


if __name__ == '__main__':
    x = 1